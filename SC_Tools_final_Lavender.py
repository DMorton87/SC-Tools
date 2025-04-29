import tkinter as tk
from tkinter import filedialog, messagebox, simpledialog, ttk
import os
import logging
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx import Presentation
from docx import Document
from docx.enum.text import WD_COLOR_INDEX
from docx.shared import Pt


# Configure logging
logging.basicConfig(
    filename="application.log",
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s"
)

# Function to show the best practices message
def show_best_practices():
    def on_close():
        if dont_show_again_var.get():
            with open("settings.txt", "w") as f:
                f.write("hide_highlight_message=True")
        best_practices_window.destroy()

    best_practices_window = tk.Toplevel()
    best_practices_window.title("Best Practices")
    best_practices_window.geometry("450x400")
    
    message = (
        "Hello, Wonderful Person,\n\n"
        "*NOTE: This program only extracts text that has been highlighted in GREEN.\n\n"
        "To get the best out of this function, please first go through your file and remove all extra formatting features, such as bolded text, italics, underlines, or variant font coloring (such as blue for FOP or purple for WBR). Thank you!"
    )

    tk.Label(best_practices_window, text=message, wraplength=400, justify="left").pack(pady=15, padx=15)

    dont_show_again_var = tk.BooleanVar()
    tk.Checkbutton(best_practices_window, text="Don't show these messages again", variable=dont_show_again_var).pack(side="right", padx=15, pady=15)

    tk.Button(best_practices_window, text="Got it!", command=on_close).pack(pady=10)
    best_practices_window.wait_window()

def show_best_practices2():
    def on_close():
        if dont_show_again_var.get():
            with open("settings.txt", "w") as f:
                f.write("hide_highlight_message=True")
        best_practices_window.destroy()

    best_practices_window = tk.Toplevel()
    best_practices_window.title("Best Practices")
    best_practices_window.geometry("450x600")
    
    message = (
        "Hello, Wonderful Person,\n\n"
        "This function will convert a properly formatted .docx file containing the and SC teachers teaching points into a .pptx that will be used on teaching day. When you run this function, you will be given the option to choose the font size along with the font color and background fill color (defaults are 42, white text, black background).\n"
        "  For the most error-free experience, make sure each teaching point entry is formatted in the following way:\n\n"
        "[Word or Phrase]\n"
        "Definition: [definition]\n"
        "Example sentence: [example sentence]\n"
    )

    tk.Label(best_practices_window, text=message, wraplength=400, justify="left").pack(pady=15, padx=15)

    dont_show_again_var = tk.BooleanVar()
    tk.Checkbutton(best_practices_window, text="Don't show these messages again", variable=dont_show_again_var).pack(side="right", padx=15, pady=15)

    tk.Button(best_practices_window, text="Got it!", command=on_close).pack(pady=10)
    best_practices_window.wait_window()

def show_best_practices3():
    def on_close():
        if dont_show_again_var.get():
            with open("settings.txt", "w") as f:
                f.write("hide_highlight_message=True")
        best_practices_window.destroy()

    best_practices_window = tk.Toplevel()
    best_practices_window.title("Best Practices")
    best_practices_window.geometry("450x300")
    
    message = (
        "Hello, Wonderful Person,\n\n"
        "Use this function to edit the formatting of an existing PowerPoint if you need to tweak the font size or color after it's been generated."
    )

    tk.Label(best_practices_window, text=message, wraplength=400, justify="left").pack(pady=15, padx=15)

    dont_show_again_var = tk.BooleanVar()
    tk.Checkbutton(best_practices_window, text="Don't show these messages again", variable=dont_show_again_var).pack(side="right", padx=15, pady=15)

    tk.Button(best_practices_window, text="Got it!", command=on_close).pack(pady=10)
    best_practices_window.wait_window()

def show_best_practices4():
    def on_close():
        if dont_show_again_var.get():
            with open("settings.txt", "w") as f:
                f.write("hide_highlight_message=True")
        best_practices_window.destroy()

    best_practices_window = tk.Toplevel()
    best_practices_window.title("Best Practices")
    best_practices_window.geometry("450x750")
    
    message = (
        "Hello, Wonderful Person,\n\n"
        "*Please check the 'SC Tools Readme' before using this function, it's rather tempermental.*\n\n" 
        "This function allows SC teachers generate a 'subtitle' .docx file from the document that contains their teaching points. To ensure the most error-free experience, please format the original teaching file in the following way:\n\n"
        "I: (Index and title information)\n"
        "O: (Opening text. If you have an especially long opening, it's fine. Just don't insert any paragraph breaks)\n"
        "T1:\n" 
        "(Lesson text)\n"
        "T2:\n" 
        "(Lesson text)\n"
        "T3:\n" 
        "(Lesson text)\n"
        "(Then O: and T1: T2: and T3: as before for day two)\n"
        "C: (Closing)\n"
    )

    tk.Label(best_practices_window, text=message, wraplength=400, justify="left").pack(pady=15, padx=15)

    dont_show_again_var = tk.BooleanVar()
    tk.Checkbutton(best_practices_window, text="Don't show these messages again", variable=dont_show_again_var).pack(side="right", padx=15, pady=15)

    tk.Button(best_practices_window, text="Got it!", command=on_close).pack(pady=10)
    best_practices_window.wait_window()

def show_best_practices5():
    def on_close():
        if dont_show_again_var.get():
            with open("settings.txt", "w") as f:
                f.write("hide_highlight_message=True")
        best_practices_window.destroy()

    best_practices_window = tk.Toplevel()
    best_practices_window.title("Best Practices")
    best_practices_window.geometry("450x400")
    
    message = (
        "Hello, Wonderful Person,\n\n"
        "This powerful function will search through a parent folder and all of its subfolders, searching for .docx files and compiling all of the text into a single plaintext file. This can process hundreds of files at a time, and is designed for gathering LLM training data (ChatGPT, Gemini, Copilot etc)."
    )

    tk.Label(best_practices_window, text=message, wraplength=400, justify="left").pack(pady=15, padx=15)

    dont_show_again_var = tk.BooleanVar()
    tk.Checkbutton(best_practices_window, text="Don't show these messages again", variable=dont_show_again_var).pack(side="right", padx=15, pady=15)

    tk.Button(best_practices_window, text="Got it!", command=on_close).pack(pady=10)
    best_practices_window.wait_window()

def show_best_practices6():
    def on_close():
        if dont_show_again_var.get():
            with open("settings.txt", "w") as f:
                f.write("hide_highlight_message=True")
        best_practices_window.destroy()

    best_practices_window = tk.Toplevel()
    best_practices_window.title("Best Practices")
    best_practices_window.geometry("450x300")
    
    message = (
        "Hello, Wonderful Person,\n\n"
        "This function performs the same task as Recursive DOCX to Plaintext, but it is specially designed to deal with files that contain embedded tables (indexes, production charts, etc)."
    )

    tk.Label(best_practices_window, text=message, wraplength=400, justify="left").pack(pady=15, padx=15)

    dont_show_again_var = tk.BooleanVar()
    tk.Checkbutton(best_practices_window, text="Don't show these messages again", variable=dont_show_again_var).pack(side="right", padx=15, pady=15)

    tk.Button(best_practices_window, text="Got it!", command=on_close).pack(pady=10)
    best_practices_window.wait_window()

def extract_highlighted_text():
    try:
        if not os.path.exists("settings.txt") or "hide_highlight_message=True" not in open("settings.txt").read():         
            show_best_practices()


        file_path = filedialog.askopenfilename(
            title="Select a .docx file",
            filetypes=[("Word Document", "*.docx")]
        )
        if not file_path:
            return

        output_folder = filedialog.askdirectory(title="Select an output folder")
        if not output_folder:
            return
        
        # Extract the base name of the file (e.g., "example" from "example.docx")
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        
        # Create the default output file name
        default_output_name = f"{base_name} Extracted Text.txt"
            
        output_file = filedialog.asksaveasfilename(
            title="What shall we call it?",
            defaultextension=".txt",
            filetypes=[("Text File", "*.txt")],
            initialfile=default_output_name         

        )
        if not output_file:
            return

        doc = Document(file_path)
        highlighted_text = []

        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                if run.font.highlight_color and run.font.highlight_color == 4:
                    highlighted_text.append(run.text.strip())

        with open(output_file, 'w', encoding='utf-8') as f:
            f.write("\n".join(highlighted_text))

        if not highlighted_text:
            messagebox.showinfo("No Highlights", "No highlighted text was found in the document.")
        else:
            logging.info(f"Extracted highlighted text from {file_path}")
        open_file = messagebox.askyesno("Success~!", f"Generated Extracted highlighted text file saved to:\n{output_file}\n\nWould you like to open it now?")
        
        if open_file:
            os.startfile(output_file)  # Opens the file using the default application
    except Exception as e:
        logging.error(f"Error extracting highlighted text: {e}")
        messagebox.showerror("Error", f"An error occurred: {e}")


def generate_powerpoint():
    try:
        if not os.path.exists("settings.txt") or "hide_highlight_message=True" not in open("settings.txt").read():         
            show_best_practices2()
        
        file_path = filedialog.askopenfilename(
            title="Select your teaching file",
            filetypes=[("Word Document", "*.docx")]
        )
        if not file_path:
            return

        output_folder = filedialog.askdirectory(title="Where do you want to save your file?")
        if not output_folder:
            return
            
        # Extract the base name of the file (e.g., "example" from "example.docx")
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        
        # Create the default output file name
        default_output_name = f"{base_name} Teaching Slides.pptx"
            
        output_file = filedialog.asksaveasfilename(
            title="What shall we call it?",
            defaultextension=".pptx",
            filetypes=[("PowerPoint Presentation", "*.pptx")],
            initialfile=default_output_name
        )
        if not output_file:
            return
        # User settings for formatting
        
        # Create a temporary root window
        root = tk.Tk()
        root.withdraw()

        # Force it to stay on top
        root.lift()
        root.attributes('-topmost', True)
        

        # Get user preferences
        body_font_size = simpledialog.askinteger("Body Font Size", "Enter body font size (default: 42)", initialvalue=42, parent=root)
        font_color = tuple(map(int, simpledialog.askstring("Font Color", "Enter font color as RGB (default: 255,255,255)", initialvalue="255,255,255", parent=root).split(",")))
        background_color = tuple(map(int, simpledialog.askstring("Background Color", "Enter background color as RGB (default: 0,0,0)", initialvalue="0,0,0", parent=root).split(",")))

        root.destroy()
        
        doc = Document(file_path)
        presentation = Presentation()

        current_title, current_definition, current_example = None, None, None

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            if current_title is None:
                current_title = text
            elif current_definition is None:
                current_definition = text.lower().replace("definition: ", "")
            elif current_example is None:
                current_example = text.lower().replace("example sentence: ", "")

                # Get user preferences
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(*background_color)


                # Add content box
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(5))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                content_frame.text = f"\n{current_title}\nD: {current_definition}\n\nE: {current_example}"
                for paragraph in content_frame.paragraphs:
                    paragraph.font.size = Pt(body_font_size)
                    paragraph.font.color.rgb = RGBColor(*font_color)

                # Reset for the next entry
                current_title, current_definition, current_example = None, None, None

        presentation.save(output_file)

        logging.info(f"Generated PowerPoint from {file_path}")
        
        open_file = messagebox.askyesno("Success~!", f"Generated PowerPoint saved to:\n{output_file}\n\nWould you like to open it now?")
        
        if open_file:
            os.startfile(output_file)  # Opens the file using the default application
    except Exception as e:
        logging.error(f"Error generating PowerPoint: {e}")
        messagebox.showerror("Error", f"An error occurred: {e}")

def format_powerpoint():
    try:
        if not os.path.exists("settings.txt") or "hide_highlight_message=True" not in open("settings.txt").read():         
            show_best_practices3()		
       
        file_path = filedialog.askopenfilename(
            title="Select a PowerPoint file",
            filetypes=[("PowerPoint Presentation", "*.pptx")]
        )
        if not file_path:
            return

        # Extract the base name of the file (e.g., "example" from "example.docx")
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        
        # Create the default output file name
        default_output_name = f"Formatted {base_name}.pptx"
            
        output_file = filedialog.asksaveasfilename(
            title="What shall we call it",
            defaultextension=".pptx",
            filetypes=[("PowerPoint Presentation", "*.pptx")],
            initialfile=default_output_name

        )
        if not output_file:
            return
            
        # Create a temporary root window
        root = tk.Tk()
        root.withdraw()

        # Force it to stay on top
        root.lift()
        root.attributes('-topmost', True)
        

        # Get user preferences
        body_font_size = simpledialog.askinteger("Body Font Size", "Enter body font size (default: 42)", initialvalue=42, parent=root)
        font_color = tuple(map(int, simpledialog.askstring("Font Color", "Enter font color as RGB (default: 255,255,255)", initialvalue="255,255,255", parent=root).split(",")))
        background_color = tuple(map(int, simpledialog.askstring("Background Color", "Enter background color as RGB (default: 0,0,0)", initialvalue="0,0,0", parent=root).split(",")))

        root.destroy()

        prs = Presentation(file_path)

        # Iterate through slides and update formatting
        for slide in prs.slides:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*background_color)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(body_font_size)
                            run.font.color.rgb = RGBColor(*font_color)

        prs.save(output_file)
        logging.info(f"Formatted PowerPoint saved as {output_file}")
        
        open_file = messagebox.askyesno("Success~!", f"Formatted PowerPoint saved to:\n{output_file}\n\nWould you like to open it now?")
        
        if open_file:
            os.startfile(output_file)  # Opens the file using the default application
    except Exception as e:
        logging.error(f"Error formatting PowerPoint: {e}")
        messagebox.showerror("Error", f"An error occurred: {e}")

def create_subtitle_file():
    try:
        if not os.path.exists("settings.txt") or "hide_highlight_message=True" not in open("settings.txt").read():         
            show_best_practices4()	
            
        file_path = filedialog.askopenfilename(
            title="Select a .docx file",
            filetypes=[("Word Document", "*.docx")]
        )
        if not file_path:
            return
            
        # Extract the base name of the file (e.g., "example" from "example.docx")
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        
        # Create the default output file name
        default_output_name = f"{base_name} Subtitle File.txt"
            
        output_file = filedialog.asksaveasfilename(
            title="What shall we call it?",
            defaultextension=".txt",
            filetypes=[("Text File", "*.txt")],
            initialfile=default_output_name

        )
        
        if not output_file:
            return

        doc = Document(file_path)
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(f"Subtitle file for {base_name}\n\n")

            section = None
            highlighted_text = []
            section_text = []

            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text.lower().startswith("i:"):
                    if section == "I:":
                        f.write("\n")
                    section = "I:"
                    f.write("\n\nIndex:\n")
                    f.write(text[2:].strip() + "\n")            
                elif text.lower().startswith("o:"):
                    section = "O:"
                    f.write("\nOpening:\n")
                    f.write(text[2:].strip() + "\n")
                elif text.lower().startswith("t1:"):
                    section = "T1:"
                    f.write("\nTeaching One:\n")
                    highlighted_text = []
                elif text.lower().startswith("t2:"):
                    if section == "T1:":
                        f.write("\n".join(highlighted_text) + "\n")
                    section = "T2:"
                    f.write("\nTeaching Two:\n")
                    highlighted_text = []
                elif text.lower().startswith("t3:"):
                    if section == "T2:":
                        f.write("\n".join(highlighted_text) + "\n")
                    section = "T3:"
                    f.write("\nTeaching Three:\n")
                    highlighted_text = []
                elif text.lower().startswith("c:"):
                    if section == "T3:":
                        f.write("\n".join(highlighted_text) + "\n")
                    section = "C:"
                    f.write("\nClosing:\n")
                    f.write(text[2:].strip() + "\n")
                else:
                    if section in ["T1:", "T2:", "T3:"]:
                        for run in paragraph.runs:
                            if run.font.highlight_color:
                                highlighted_text.append(run.text.strip())

            if section in ["T1:", "T2:", "T3:"]:
                f.write("\n".join(highlighted_text) + "\n")

        open_file = messagebox.askyesno("Success~!", f"Generated Subtitle file and saved it to:\n{output_file}\n\nWould you like to open it now?")
        
        if open_file:
            os.startfile(output_file)  # Opens the file using the default application
    except Exception as e:
        logging.error(f"Error creating subtitle file: {e}")
        messagebox.showerror("Error", f"An error occurred: {e}")



def recursive_docx_to_plaintext():
    try:
        if not os.path.exists("settings.txt") or "hide_highlight_message=True" not in open("settings.txt").read():         
            show_best_practices5()
            	          
        parent_folder = filedialog.askdirectory(title="Select a parent folder")
        if not parent_folder:
            return

        output_file = filedialog.asksaveasfilename(
            title="Save extracted text as",
            defaultextension=".txt",
            filetypes=[("Text File", "*.txt")],
            initialfile="extracted_text.txt"
        )
        if not output_file:
            return

        # Create a progress window
        progress_window = tk.Toplevel()
        progress_window.title("Processing...")
        progress_window.geometry("400x100")
        tk.Label(progress_window, text="Extracting text and tables...").pack(pady=10)

        progress = ttk.Progressbar(progress_window, orient="horizontal", length=300, mode="determinate")
        progress.pack(pady=10)

        # Count total files
        total_files = sum(len(files) for _, _, files in os.walk(parent_folder) if any(f.endswith(".docx") for f in files))
        progress["maximum"] = total_files

        current_file = 0

        with open(output_file, "w", encoding="utf-8") as output:
            for root, _, files in os.walk(parent_folder):
                for file_name in files:
                    if file_name.endswith(".docx") and not file_name.startswith("~$"):
                        doc_path = os.path.join(root, file_name)
                        try:
                            doc = Document(doc_path)
                            output.write(f"=== {file_name} ===\n")
                            output.write("\n".join([paragraph.text for paragraph in doc.paragraphs]) + "\n\n")

                            current_file += 1
                            progress["value"] = current_file
                            progress_window.update_idletasks()
                        except Exception as e:
                            logging.error(f"Error processing file {doc_path}: {e}")
                            continue

        logging.info(f"Extraction completed. Output saved to {output_file}")
        messagebox.showinfo("Success", f"Extraction completed! Saved to {output_file}")
    except Exception as e:
        logging.error(f"Error during extraction: {e}")
        messagebox.showerror("Error", f"An error occurred: {e}")
    finally:
        progress_window.destroy()


def recursive_index_to_plaintext():
    try:
        if not os.path.exists("settings.txt") or "hide_highlight_message=True" not in open("settings.txt").read():         
            show_best_practices6()	
                
        parent_folder = filedialog.askdirectory(title="Select a parent folder")
        if not parent_folder:
            return

        output_file = filedialog.asksaveasfilename(
            title="Save extracted text as",
            defaultextension=".txt",
            filetypes=[("Text File", "*.txt")],
            initialfile="extracted_text.txt"
        )
        if not output_file:
            return

        # Create a progress window
        progress_window = tk.Toplevel()
        progress_window.title("Processing...")
        progress_window.geometry("400x100")
        tk.Label(progress_window, text="Extracting text and tables...").pack(pady=10)

        progress = ttk.Progressbar(progress_window, orient="horizontal", length=300, mode="determinate")
        progress.pack(pady=10)

        # Count total files
        total_files = sum(len(files) for _, _, files in os.walk(parent_folder) if any(f.endswith(".docx") for f in files))
        progress["maximum"] = total_files

        current_file = 0

        with open(output_file, "w", encoding="utf-8") as output:
            for root, _, files in os.walk(parent_folder):
                for file_name in files:
                    if file_name.endswith(".docx") and not file_name.startswith("~$"):
                        doc_path = os.path.join(root, file_name)
                        try:
                            doc = Document(doc_path)
                            output.write(f"=== {file_name} ===\n")
                            output.write("\n".join([paragraph.text for paragraph in doc.paragraphs]) + "\n\n")

                            for table in doc.tables:
                                output.write("Table:\n")
                                for row in table.rows:
                                    row_data = [cell.text.strip() for cell in row.cells]
                                    output.write("\t".join(row_data) + "\n")
                                output.write("\n")

                            current_file += 1
                            progress["value"] = current_file
                            progress_window.update_idletasks()
                        except Exception as e:
                            logging.error(f"Error processing file {doc_path}: {e}")
                            continue

        logging.info(f"Extraction completed. Output saved to {output_file}")
        messagebox.showinfo("Success", f"Extraction completed! Saved to {output_file}")
    except Exception as e:
        logging.error(f"Error during extraction: {e}")
        messagebox.showerror("Error", f"An error occurred: {e}")
    finally:
        progress_window.destroy()

# Main GUI application
def main():
    root = tk.Tk()
    root.title("~ SC Tools Lavender ~")
    root.geometry("400x500")
    root.configure(bg="lavender blush")

    tk.Label(root, text="**~~ Welcome to SC Tools! ~~**", font=("Georgia", 18), bg="lavender blush").pack(pady=10)

 
    tk.Button(
        root, text="Extract Highlighted Text", command=extract_highlighted_text, width=25,
        bg="snow"
    ).pack(pady=10)

    tk.Button(
        root, text="Generate PowerPoint", command=generate_powerpoint, width=25,
        bg="lavender"
    ).pack(pady=10)
    
    tk.Button(
        root, text="Format PowerPoint", command=format_powerpoint, width=25,
        bg="snow"
    ).pack(pady=10)
    
    tk.Button(
        root, text="Generate Subtitle File", command=create_subtitle_file, width=25,
        bg="lavender" 
    ).pack(pady=10)

    tk.Button(
        root, text="Recursive .docx to Plaintext", command=recursive_docx_to_plaintext, width=25,
        bg="snow"
    ).pack(pady=10)
   
    tk.Button(
        root, text="Recursive Index to Plaintext", command=recursive_index_to_plaintext, width=25,
        bg="lavender"
    ).pack(pady=10)





    root.mainloop()

if __name__ == "__main__":
    main()
